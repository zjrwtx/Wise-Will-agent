"""
PowerPoint presentation builder using python-pptx.

This module provides functionality to build PowerPoint presentations
from slide content using various layouts and styles.

Example:
    >>> builder = PptBuilder()
    >>> ppt_path = builder.build(slides, "output.pptx", "My Presentation")
    >>> print(f"PPT saved to: {ppt_path}")
"""

from __future__ import annotations

import logging
from pathlib import Path
from typing import Optional

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

from .models import SlideContent

logger = logging.getLogger(__name__)


# Color schemes for different styles
COLOR_SCHEMES = {
    "professional": {
        "primary": RGBColor(0x1A, 0x73, 0xE8),      # Google Blue
        "secondary": RGBColor(0x34, 0xA8, 0x53),    # Google Green
        "accent": RGBColor(0xFB, 0xBC, 0x04),       # Google Yellow
        "text": RGBColor(0x20, 0x21, 0x24),         # Dark gray
        "light": RGBColor(0xF8, 0xF9, 0xFA),        # Light gray
    },
    "academic": {
        "primary": RGBColor(0x00, 0x33, 0x66),      # Navy blue
        "secondary": RGBColor(0x8B, 0x00, 0x00),    # Dark red
        "accent": RGBColor(0xCC, 0x99, 0x00),       # Gold
        "text": RGBColor(0x33, 0x33, 0x33),         # Dark gray
        "light": RGBColor(0xF5, 0xF5, 0xF5),        # Off white
    },
    "creative": {
        "primary": RGBColor(0x66, 0x33, 0x99),      # Purple
        "secondary": RGBColor(0xFF, 0x66, 0x99),    # Pink
        "accent": RGBColor(0x00, 0xCC, 0xCC),       # Teal
        "text": RGBColor(0x2D, 0x2D, 0x2D),         # Charcoal
        "light": RGBColor(0xFD, 0xFD, 0xFD),        # White
    },
}


class PptBuilder:
    """Builder for PowerPoint presentations.

    This class creates PowerPoint presentations from slide content
    using python-pptx library with customizable styles and layouts.

    Attributes:
        style: Presentation style (professional, academic, creative).
        colors: Color scheme for the presentation.

    Example:
        >>> builder = PptBuilder(style="professional")
        >>> ppt_path = builder.build(slides, "output.pptx", "My PPT")
    """

    def __init__(self, style: str = "professional"):
        """Initialize the PPT builder.

        Args:
            style: Presentation style. One of: professional, academic,
                   creative.

        Example:
            >>> builder = PptBuilder(style="academic")
        """
        self.style = style
        self.colors = COLOR_SCHEMES.get(style, COLOR_SCHEMES["professional"])

    def build(
        self,
        slides: list[SlideContent],
        output_path: str | Path,
        title: str = "Presentation",
        author: str = "AI Generated",
    ) -> Path:
        """Build a PowerPoint presentation from slide content.

        Args:
            slides: List of SlideContent objects.
            output_path: Path for the output .pptx file.
            title: Presentation title.
            author: Author name for metadata.

        Returns:
            Path to the generated .pptx file.

        Example:
            >>> builder = PptBuilder()
            >>> slides = [SlideContent(index=0, title="Hello", layout="title")]
            >>> path = builder.build(slides, "output.pptx", "My PPT")
            >>> print(f"Created: {path}")
        """
        output_path = Path(output_path)
        output_path.parent.mkdir(parents=True, exist_ok=True)

        logger.info(f"Building PPT: {output_path} ({len(slides)} slides)")

        # Create presentation
        prs = Presentation()
        prs.slide_width = Inches(13.333)  # 16:9 aspect ratio
        prs.slide_height = Inches(7.5)

        # Set metadata
        prs.core_properties.title = title
        prs.core_properties.author = author

        # Build each slide
        for slide_content in slides:
            self._add_slide(prs, slide_content)

        # Save presentation
        prs.save(str(output_path))
        logger.info(f"PPT saved: {output_path}")

        return output_path

    def _add_slide(self, prs: Presentation, content: SlideContent) -> None:
        """Add a slide to the presentation.

        Args:
            prs: Presentation object.
            content: Slide content to add.
        """
        layout = content.layout.lower()

        if layout == "title":
            self._add_title_slide(prs, content)
        elif layout == "bullets":
            self._add_bullets_slide(prs, content)
        elif layout == "two_column":
            self._add_two_column_slide(prs, content)
        elif layout == "summary":
            self._add_summary_slide(prs, content)
        else:
            self._add_content_slide(prs, content)

    def _add_title_slide(
        self,
        prs: Presentation,
        content: SlideContent,
    ) -> None:
        """Add a title slide.

        Args:
            prs: Presentation object.
            content: Slide content.
        """
        blank_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(blank_layout)

        # Background shape
        bg_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            prs.slide_width, prs.slide_height
        )
        bg_shape.fill.solid()
        bg_shape.fill.fore_color.rgb = self.colors["primary"]
        bg_shape.line.fill.background()

        # Title
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(2.5),
            Inches(11.333), Inches(1.5)
        )
        title_frame = title_box.text_frame
        title_frame.paragraphs[0].text = content.title
        title_frame.paragraphs[0].font.size = Pt(48)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Subtitle (from content or notes)
        subtitle_text = content.content or content.notes
        if subtitle_text:
            subtitle_box = slide.shapes.add_textbox(
                Inches(1), Inches(4.2),
                Inches(11.333), Inches(1)
            )
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.paragraphs[0].text = subtitle_text
            subtitle_frame.paragraphs[0].font.size = Pt(24)
            subtitle_frame.paragraphs[0].font.color.rgb = RGBColor(
                0xFF, 0xFF, 0xFF
            )
            subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    def _add_content_slide(
        self,
        prs: Presentation,
        content: SlideContent,
    ) -> None:
        """Add a content slide with title and body text.

        Args:
            prs: Presentation object.
            content: Slide content.
        """
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)

        # Title bar
        title_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            prs.slide_width, Inches(1.2)
        )
        title_bar.fill.solid()
        title_bar.fill.fore_color.rgb = self.colors["primary"]
        title_bar.line.fill.background()

        # Title text
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3),
            Inches(12.333), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_frame.paragraphs[0].text = content.title
        title_frame.paragraphs[0].font.size = Pt(32)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

        # Content
        content_text = content.content
        if not content_text and content.bullet_points:
            content_text = "\n".join(f"• {bp}" for bp in content.bullet_points)

        if content_text:
            content_box = slide.shapes.add_textbox(
                Inches(0.75), Inches(1.6),
                Inches(11.833), Inches(5.5)
            )
            content_frame = content_box.text_frame
            content_frame.word_wrap = True

            # Split into paragraphs
            paragraphs = content_text.split("\n")
            for i, para_text in enumerate(paragraphs):
                if i == 0:
                    p = content_frame.paragraphs[0]
                else:
                    p = content_frame.add_paragraph()
                p.text = para_text
                p.font.size = Pt(20)
                p.font.color.rgb = self.colors["text"]
                p.space_after = Pt(12)

    def _add_bullets_slide(
        self,
        prs: Presentation,
        content: SlideContent,
    ) -> None:
        """Add a slide with bullet points.

        Args:
            prs: Presentation object.
            content: Slide content.
        """
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)

        # Title bar
        title_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            prs.slide_width, Inches(1.2)
        )
        title_bar.fill.solid()
        title_bar.fill.fore_color.rgb = self.colors["primary"]
        title_bar.line.fill.background()

        # Title text
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3),
            Inches(12.333), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_frame.paragraphs[0].text = content.title
        title_frame.paragraphs[0].font.size = Pt(32)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

        # Bullet points
        if content.bullet_points:
            bullets_box = slide.shapes.add_textbox(
                Inches(0.75), Inches(1.6),
                Inches(11.833), Inches(5.5)
            )
            bullets_frame = bullets_box.text_frame
            bullets_frame.word_wrap = True

            for i, bullet in enumerate(content.bullet_points):
                if i == 0:
                    p = bullets_frame.paragraphs[0]
                else:
                    p = bullets_frame.add_paragraph()

                # Add bullet marker
                p.text = f"●  {bullet}"
                p.font.size = Pt(22)
                p.font.color.rgb = self.colors["text"]
                p.space_before = Pt(8)
                p.space_after = Pt(8)
                p.level = 0

    def _add_two_column_slide(
        self,
        prs: Presentation,
        content: SlideContent,
    ) -> None:
        """Add a two-column slide.

        Args:
            prs: Presentation object.
            content: Slide content.
        """
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)

        # Title bar
        title_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            prs.slide_width, Inches(1.2)
        )
        title_bar.fill.solid()
        title_bar.fill.fore_color.rgb = self.colors["primary"]
        title_bar.line.fill.background()

        # Title text
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3),
            Inches(12.333), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_frame.paragraphs[0].text = content.title
        title_frame.paragraphs[0].font.size = Pt(32)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

        # Left column
        left_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.6),
            Inches(5.916), Inches(5.5)
        )
        left_frame = left_box.text_frame
        left_frame.word_wrap = True
        left_text = content.left_content or (
            "\n".join(content.bullet_points[:len(content.bullet_points)//2])
            if content.bullet_points else ""
        )
        if left_text:
            left_frame.paragraphs[0].text = left_text
            left_frame.paragraphs[0].font.size = Pt(18)
            left_frame.paragraphs[0].font.color.rgb = self.colors["text"]

        # Right column
        right_box = slide.shapes.add_textbox(
            Inches(6.916), Inches(1.6),
            Inches(5.916), Inches(5.5)
        )
        right_frame = right_box.text_frame
        right_frame.word_wrap = True
        right_text = content.right_content or (
            "\n".join(content.bullet_points[len(content.bullet_points)//2:])
            if content.bullet_points else ""
        )
        if right_text:
            right_frame.paragraphs[0].text = right_text
            right_frame.paragraphs[0].font.size = Pt(18)
            right_frame.paragraphs[0].font.color.rgb = self.colors["text"]

        # Divider line
        divider = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(6.666), Inches(1.6),
            Inches(0.02), Inches(5.5)
        )
        divider.fill.solid()
        divider.fill.fore_color.rgb = self.colors["light"]
        divider.line.fill.background()

    def _add_summary_slide(
        self,
        prs: Presentation,
        content: SlideContent,
    ) -> None:
        """Add a summary/conclusion slide.

        Args:
            prs: Presentation object.
            content: Slide content.
        """
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)

        # Background
        bg_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            prs.slide_width, prs.slide_height
        )
        bg_shape.fill.solid()
        bg_shape.fill.fore_color.rgb = self.colors["secondary"]
        bg_shape.line.fill.background()

        # Title
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(1),
            Inches(11.333), Inches(1.2)
        )
        title_frame = title_box.text_frame
        title_frame.paragraphs[0].text = content.title
        title_frame.paragraphs[0].font.size = Pt(40)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Summary points
        if content.bullet_points:
            summary_box = slide.shapes.add_textbox(
                Inches(1.5), Inches(2.5),
                Inches(10.333), Inches(4.5)
            )
            summary_frame = summary_box.text_frame
            summary_frame.word_wrap = True

            for i, point in enumerate(content.bullet_points):
                if i == 0:
                    p = summary_frame.paragraphs[0]
                else:
                    p = summary_frame.add_paragraph()
                p.text = f"✓  {point}"
                p.font.size = Pt(24)
                p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                p.space_after = Pt(16)
                p.alignment = PP_ALIGN.LEFT

        # Content text
        elif content.content:
            content_box = slide.shapes.add_textbox(
                Inches(1.5), Inches(2.5),
                Inches(10.333), Inches(4.5)
            )
            content_frame = content_box.text_frame
            content_frame.word_wrap = True
            content_frame.paragraphs[0].text = content.content
            content_frame.paragraphs[0].font.size = Pt(22)
            content_frame.paragraphs[0].font.color.rgb = RGBColor(
                0xFF, 0xFF, 0xFF
            )
            content_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
